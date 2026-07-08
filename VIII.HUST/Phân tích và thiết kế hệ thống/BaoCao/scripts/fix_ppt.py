from pptx import Presentation

prs = Presentation(r"c:\Dev\Learning\VIII.HUST\Phân tích và thiết kế hệ thống\request\Bai Tap Lon\Nhóm 5\Nhom05_Slide_Thuyet_Trinh.pptx")

print(f"Total slides: {len(prs.slides)}")
for i, slide in enumerate(prs.slides):
    print(f"\nSlide {i+1}:")
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            print(f"  - {shape.text[:100]}")
